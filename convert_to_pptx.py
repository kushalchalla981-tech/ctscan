"""Convert presentation.html to PPTX with dark theme matching the HTML style."""
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0A, 0x0A, 0x0F)
WHITE = RGBColor(0xE8, 0xE8, 0xED)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)
DIM = RGBColor(0x6A, 0x6A, 0x7A)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
PURPLE = RGBColor(0xA5, 0xB4, 0xFC)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED = RGBColor(0xF8, 0x71, 0x71)
CARD_BG = RGBColor(0x14, 0x14, 0x1A)
CARD_BORDER = RGBColor(0x22, 0x22, 0x2E)

def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_txt(slide, text, left, top, width, height, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name='Calibri'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tb

def add_para(text_frame, text, font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name='Calibri', spacing_before=0, spacing_after=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = Pt(spacing_before)
    p.space_after = Pt(spacing_after)
    return p

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(0.5)
    return shape

def add_run(p, text, font_size=14, bold=False, color=WHITE, font_name='Calibri'):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return run

# ═══ Read HTML ═══
with open('presentation.html', 'r', encoding='utf-8') as f:
    html = f.read()

# Split by slide sections
slides_html = re.findall(r'<section class="slide"[^>]*>(.*?)</section>', html, re.DOTALL)

def strip_tags(s):
    return re.sub(r'<[^>]+>', '', s).strip()

def get_inner_html(match):
    """Get the inner text of a tag, stripping all HTML tags"""
    return strip_tags(match.group(0))

for idx, slide_html in enumerate(slides_html):
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)

    # Extract slide number
    sn_match = re.search(r'<div class="slide-number">\s*(\d+)\s*/\s*(\d+)', slide_html)
    slide_num = sn_match.group(1) if sn_match else str(idx + 1)
    total_slides = sn_match.group(2) if sn_match else str(len(slides_html))

    # Add slide number bottom-right
    add_txt(slide, f"{slide_num} / {total_slides}", 11.5, 7.0, 1.5, 0.4,
            font_size=11, color=DIM, align=PP_ALIGN.RIGHT)

    # ── SLIDE 1: TITLE ──
    if idx == 0:
        # Tags
        tags_match = re.findall(r'<span class="tag[^"]*"[^>]*>([^<]+)</span>', slide_html)
        if tags_match:
            tags_text = "  ".join(tags_match)
            add_txt(slide, tags_text, 2.5, 1.5, 8, 0.5, font_size=14, color=AMBER, align=PP_ALIGN.CENTER)

        # Title
        t_match = re.search(r'<h1[^>]*>(.*?)</h1>', slide_html, re.DOTALL)
        if t_match:
            title = strip_tags(t_match.group(1)).replace('<br>', '\n')
            add_txt(slide, title, 1.5, 2.2, 10.3, 1.5, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Subtitle
        sub_match = re.search(r'<p class="subtitle[^"]*"[^>]*>(.*?)</p>', slide_html, re.DOTALL)
        if sub_match:
            subtitle = strip_tags(sub_match.group(1)).replace('<br>', '\n')
            add_txt(slide, subtitle, 2, 3.8, 9.3, 1.0, font_size=20, color=GRAY, align=PP_ALIGN.CENTER)

        # Stats
        stats = re.findall(r'<span class="stat-value[^"]*"[^>]*>([^<]+)</span>\s*<div class="stat-label">([^<]+)</div>', slide_html)
        if stats:
            x_start = 2.0
            for i, (val, label) in enumerate(stats):
                ax = x_start + i * 2.8
                add_txt(slide, val, ax, 5.0, 2.5, 0.6, font_size=36, bold=True,
                        color=[PURPLE, GREEN, AMBER, RED][i % 4], align=PP_ALIGN.CENTER)
                add_txt(slide, label, ax, 5.6, 2.5, 0.4, font_size=12, color=DIM, align=PP_ALIGN.CENTER)
        continue

    # ── Extract h3 and h2 for all slides ──
    h3_match = re.search(r'<h3>(.*?)</h3>', slide_html)
    h2_match = re.search(r'<h2>(.*?)</h2>', slide_html, re.DOTALL)

    y_cursor = 0.4
    # h3 label
    if h3_match:
        add_txt(slide, strip_tags(h3_match.group(1)).upper(), 0.8, y_cursor, 10, 0.4,
                font_size=13, bold=True, color=AMBER)
        y_cursor += 0.45

    # h2 title
    if h2_match:
        h2_text = strip_tags(h2_match.group(1))
        add_txt(slide, h2_text, 0.8, y_cursor, 11, 0.6, font_size=28, bold=True, color=WHITE)
        y_cursor += 0.65

    # Main paragraph (the first <p> after h2 that isn't inside a card)
    p_matches = re.finditer(r'<p[^>]*>(.*?)</p>', slide_html, re.DOTALL)
    for pm in p_matches:
        p_text = strip_tags(pm.group(1))
        # Only take non-empty, non-card paragraphs at the top level
        if len(p_text) > 30 and 'Measured' in p_text or 'Real CT' in p_text or 'Every formula' in p_text:
            add_txt(slide, p_text, 0.8, y_cursor, 11.5, 0.6, font_size=14, color=GRAY)
            y_cursor += 0.55
            break

    # ── SLIDE 2: BIG PICTURE ──
    if idx == 1:
        # Math block
        mb = re.search(r'<div class="math-block[^>]*">(.*?)</div>', slide_html, re.DOTALL)
        if mb:
            mb_text = strip_tags(mb.group(1))
            # Card for the math block
            add_card(slide, 2.5, y_cursor + 0.1, 7.5, 0.7)
            add_txt(slide, mb_text, 2.5, y_cursor + 0.15, 7.5, 0.6, font_size=18, color=WHITE, align=PP_ALIGN.CENTER, font_name='Consolas')
            y_cursor += 0.9

        # Two cards: Forward + Inverse
        cards = re.findall(r'<div class="card">(.*?)</div>\s*</div>\s*</div>', slide_html, re.DOTALL)
        # Actually just find the two main cards
        card_sections = re.findall(r'<div class="card">\s*<span class="card-icon">[^<]+</span>\s*<h4>(.*?)</h4>\s*<p>(.*?)</p>', slide_html, re.DOTALL)
        if len(card_sections) >= 2:
            for i, (h4, cp) in enumerate(card_sections[:2]):
                cx = 0.8 + i * 5.8
                add_card(slide, cx, y_cursor, 5.5, 2.2)
                add_txt(slide, strip_tags(h4), cx + 0.3, y_cursor + 0.2, 5.0, 0.4, font_size=16, bold=True, color=WHITE)
                # Card text - split into lines
                card_text = strip_tags(cp)
                add_txt(slide, card_text, cx + 0.3, y_cursor + 0.7, 5.0, 1.3, font_size=12, color=GRAY)
            y_cursor += 2.5

        # Key insight
        ins_match = re.search(r'<div class="insight">\s*<p>(.*?)</p>', slide_html, re.DOTALL)
        if ins_match:
            ins_text = strip_tags(ins_match.group(1))
            add_txt(slide, ins_text, 0.8, y_cursor, 11.5, 0.6, font_size=13, color=GRAY)
        continue

    # ── SLIDE 3: FORWARD MODEL ──
    if idx == 2:
        # 3 cards for phantom, ray tracing, sinogram
        cards_data = re.findall(r'<h4>(.*?)</h4>\s*<p>(.*?)</p>', slide_html, re.DOTALL)
        if cards_data:
            for i, (h4, cp) in enumerate(cards_data):
                cx = 0.8 + i * 3.9
                add_card(slide, cx, y_cursor, 3.6, 2.0)
                add_txt(slide, strip_tags(h4), cx + 0.2, y_cursor + 0.1, 3.2, 0.4, font_size=14, bold=True, color=WHITE)
                card_text = strip_tags(cp)
                add_txt(slide, card_text, cx + 0.2, y_cursor + 0.6, 3.2, 1.2, font_size=11, color=GRAY)
                # Arrow line
                arrow_match = re.search(r'<div class="mt-4 op-70 text-sm">→\s*(.*?)</div>', slide_html, re.DOTALL)
            y_cursor += 2.3

        # Stats row
        stats = re.findall(r'<div class="stat-value[^"]*"[^>]*>([^<]+)</div>\s*<div class="stat-label">([^<]+)</div>\s*<div class="stat-desc">([^<]+)</div>', slide_html, re.DOTALL)
        if stats:
            for i, (val, label, desc) in enumerate(stats):
                cx = 0.8 + i * 3.9
                add_card(slide, cx, y_cursor, 3.6, 1.0)
                add_txt(slide, val, cx + 0.2, y_cursor + 0.05, 3.2, 0.4, font_size=24, bold=True,
                        color=[PURPLE, AMBER, GREEN][i % 3], align=PP_ALIGN.CENTER)
                add_txt(slide, label, cx + 0.2, y_cursor + 0.45, 3.2, 0.25, font_size=9, color=DIM, align=PP_ALIGN.CENTER)
                add_txt(slide, desc, cx + 0.2, y_cursor + 0.65, 3.2, 0.25, font_size=8, color=DIM, align=PP_ALIGN.CENTER)
        continue

    # ── SLIDE 4: THE MATH ──
    if idx == 3:
        # Two main cards: Least-Squares + Two Solver Strategies
        ls_match = re.search(r'<h4><span class="accent-purple">▸</span>\s*Least-Squares Solution</h4>\s*<p[^>]*>(.*?)</p>', slide_html, re.DOTALL)
        if ls_match:
            add_card(slide, 0.8, y_cursor, 5.5, 2.0)
            add_txt(slide, "Least-Squares Solution", 1.1, y_cursor + 0.1, 5.0, 0.4, font_size=15, bold=True, color=PURPLE)
            ls_text = strip_tags(ls_match.group(1))
            add_txt(slide, ls_text, 1.1, y_cursor + 0.5, 5.0, 0.4, font_size=11, color=GRAY)
            # math block
            mb = re.search(r'<div class="math-block[^>]*">(.*?)</div>', slide_html, re.DOTALL)
            if mb:
                mb_text = strip_tags(mb.group(1))
                add_txt(slide, mb_text, 1.1, y_cursor + 0.9, 5.0, 0.4, font_size=14, color=WHITE, align=PP_ALIGN.CENTER, font_name='Consolas')
            # normal equations
            ne_match = re.search(r'<span class="math-inline[^>]*">(.*?)</span>', slide_html, re.DOTALL)
            if ne_match:
                ne_text = strip_tags(ne_match.group(1))
                add_txt(slide, f"The solution satisfies the normal equations: {ne_text}", 1.1, y_cursor + 1.3, 5.0, 0.4, font_size=10, color=DIM)

        # Solver Strategies card
        add_card(slide, 6.6, y_cursor, 5.8, 2.0)
        add_txt(slide, "Two Solver Strategies", 6.9, y_cursor + 0.1, 5.3, 0.4, font_size=15, bold=True, color=GREEN)

        # Dense LU sub-card
        lu_match = re.search(r'<strong[^>]*>Dense LU</strong>\s*<p[^>]*>(.*?)</p>', slide_html, re.DOTALL)
        if lu_match:
            add_txt(slide, "Dense LU", 7.1, y_cursor + 0.55, 5.0, 0.3, font_size=13, bold=True, color=PURPLE)
            add_txt(slide, strip_tags(lu_match.group(1)), 7.1, y_cursor + 0.85, 5.0, 0.4, font_size=10, color=GRAY)

        # LSQR sub-card
        lsqr_match = re.search(r'<strong[^>]*>Sparse LSQR</strong>\s*<p[^>]*>(.*?)</p>', slide_html, re.DOTALL)
        if lsqr_match:
            add_txt(slide, "Sparse LSQR", 7.1, y_cursor + 1.25, 5.0, 0.3, font_size=13, bold=True, color=GREEN)
            add_txt(slide, strip_tags(lsqr_match.group(1)), 7.1, y_cursor + 1.55, 5.0, 0.4, font_size=10, color=GRAY)

        y_cursor += 2.3

        # Bottom 3 stat cards
        stats = re.findall(r'<div class="stat-value[^"]*"[^>]*>([^<]+)</div>\s*<div class="stat-label[^"]*"[^>]*>([^<]+)</div>\s*<div class="stat-desc">([^<]+)</div>', slide_html, re.DOTALL)
        if stats:
            for i, (val, label, desc) in enumerate(stats):
                cx = 0.8 + i * 3.9
                add_card(slide, cx, y_cursor, 3.6, 1.0)
                add_txt(slide, val, cx + 0.2, y_cursor + 0.05, 3.2, 0.35, font_size=14, bold=True,
                        color=[AMBER, GREEN, PURPLE][i % 3], align=PP_ALIGN.CENTER)
                add_txt(slide, label, cx + 0.2, y_cursor + 0.4, 3.2, 0.25, font_size=8, color=DIM, align=PP_ALIGN.CENTER)
                add_txt(slide, desc, cx + 0.2, y_cursor + 0.65, 3.2, 0.3, font_size=9, color=GRAY, align=PP_ALIGN.CENTER, font_name='Consolas')
        continue

    # ── SLIDE 5: KEY FORMULAE ──
    if idx == 4:
        # 6 formula cards in 2 columns
        col1_cards = [
            ("Forward Model", PURPLE,
             ["A · x = b", "min ‖A·x − b‖²"],
             "projector.py — builds sparse A, b\nlud_solver.py — all solver strategies"),
            ("LU Decomposition", GREEN,
             ["P·A = L·U", "L·y = P·b  ⟶  U·x = y"],
             "lud_solver.py:113–153 — dense PA=LU\nwith partial pivoting"),
            ("Tikhonov Regularization", AMBER,
             ["min ‖A·x−b‖² + λ²·‖x‖²"],
             "lud_solver.py — LSQR damp = √λ\nnoise.py — robustness sweep (λ=4.0)"),
        ]
        col2_cards = [
            ("Iterative Refinement", RED,
             ["rₖ = b − A·xₖ", "xₖ₊₁ = xₖ + A†·rₖ"],
             "lud_solver.py:156–191 — repeats until\n‖r‖/‖b‖ < 10⁻¹⁰ (max 3 iterations)"),
            ("Quality Metrics", PURPLE,
             ["RMSE = √(mean((x̂−x)²))", "PSNR = 20·log₁₀(max/RMSE)", "SSIM = (2μₓμᵧ+C₁)(2σₓᵧ+C₂) / ((μₓ²+μᵧ²+C₁)(σₓ²+σᵧ²+C₂))"],
             "metrics.py:23–76 — RMSE, PSNR, SSIM,\nrelative error, forward residual"),
            ("Other Key Formulae", RGBColor(0xC0, 0x84, 0xFC),
             ["sparsity = 1 − nnz/(m·n)", "noise = b + σ·N(0,1), σ = noise·std(b)"],
             "projector.py — sparsity calc\nnoise.py — Gaussian & Poisson"),
        ]

        for col_idx, cards in enumerate([col1_cards, col2_cards]):
            cx = 0.8 + col_idx * 5.8
            for row_idx, (title, color, formulae, source) in enumerate(cards):
                cy = y_cursor + row_idx * 1.65
                add_card(slide, cx, cy, 5.5, 1.5)
                add_txt(slide, f"▸ {title}", cx + 0.2, cy + 0.05, 5.0, 0.3, font_size=13, bold=True, color=color)
                for fi, formula in enumerate(formulae):
                    # Smaller formula sub-card
                    add_txt(slide, formula, cx + 0.2, cy + 0.35 + fi * 0.3, 5.0, 0.3,
                            font_size=11, color=WHITE, align=PP_ALIGN.CENTER, font_name='Consolas')
                add_txt(slide, source, cx + 0.2, cy + 1.05, 5.0, 0.4, font_size=8, color=DIM)

        # Insight
        add_txt(slide, "18+ unique formulae across 4 phases, each linked to its source implementation.", 0.8, 6.6, 11, 0.4, font_size=11, color=GRAY)
        continue

    # ── SLIDE 6: ARCHITECTURE ──
    if idx == 5:
        # Left column: Code block showing project structure
        add_card(slide, 0.8, y_cursor, 5.5, 3.0)

        # Extract code text preserving line structure
        code_match = re.search(r'<div class="code-block">(.*?)</div>', slide_html, re.DOTALL)
        if code_match:
            code_html = code_match.group(1)
            # Break into lines based on <br>
            code_lines = re.split(r'<br>\s*', code_html)
            code_texts = []
            for line in code_lines:
                line = re.sub(r'<span[^>]*>', '', line)
                line = re.sub(r'</span>', '', line)
                line = line.replace('&nbsp;', ' ')
                line = line.replace('&amp;', '&')
                line = line.replace('&nbsp;&nbsp;', '  ')
                line = line.replace('&nbsp;&nbsp;&nbsp;', '   ')
                line = line.replace('─', ' ')
                line = line.strip()
                code_texts.append(line)

            # Add code title
            add_txt(slide, "Project Structure", 1.1, y_cursor + 0.05, 5.0, 0.3, font_size=13, bold=True, color=WHITE)
            # Code lines with better indentation
            code_display = "\n".join(code_texts)
            add_txt(slide, code_display, 1.1, y_cursor + 0.4, 5.0, 2.4, font_size=10, color=GRAY, font_name='Consolas')

        # Right column: CLI Commands
        add_card(slide, 6.6, y_cursor, 5.8, 3.0)
        add_txt(slide, "CLI Commands", 6.9, y_cursor + 0.05, 5.3, 0.3, font_size=13, bold=True, color=WHITE)

        # Extract CLI entries from the HTML card
        cli_html = re.search(r'<div class="card"[^>]*>.*?<h4>CLI Commands</h4>(.*?)</div>', slide_html, re.DOTALL)
        if cli_html:
            cli_content = cli_html.group(1)
            cmd_entries = re.findall(
                r'<span style="color:#c8c8d0;">(.*?)</span>\s*<span class="text-xs op-70">(.*?)</span>',
                cli_content, re.DOTALL
            )
            for i, (cmd, desc) in enumerate(cmd_entries):
                cy = y_cursor + 0.45 + i * 0.5
                cmd_clean = strip_tags(cmd).strip()
                desc_clean = strip_tags(desc).strip()
                # Command in monospace
                add_txt(slide, f"$ python main.py {cmd_clean}", 7.1, cy, 5.0, 0.25,
                        font_size=11, color=AMBER, font_name='Consolas')
                # Description below
                add_txt(slide, desc_clean, 7.1, cy + 0.25, 5.0, 0.2,
                        font_size=9, color=DIM)

        y_cursor += 3.3

        # Key insight at the bottom
        ins_match = re.search(r'<div class="insight">\s*<p>(.*?)</p>', slide_html, re.DOTALL)
        if ins_match:
            ins_text = strip_tags(ins_match.group(1))
            add_txt(slide, ins_text, 0.8, y_cursor, 11.5, 0.7, font_size=11, color=GRAY)
        continue

    # ── SLIDE 7: RESULTS ──
    if idx == 6:
        # Table
        table_data = []
        rows = re.findall(r'<tr>(.*?)</tr>', slide_html, re.DOTALL)
        for row in rows[1:]:  # skip header
            cells = re.findall(r'<td[^>]*>(.*?)</td>', row, re.DOTALL)
            if cells:
                table_data.append([strip_tags(c) for c in cells])

        if table_data:
            tbl_left = 0.8
            tbl_top = y_cursor
            tbl_width = 5.5
            col_widths = [1.2, 1.2, 1.5, 1.6]
            num_rows = len(table_data) + 1
            num_cols = 4

            # Create table
            tbl_shape = slide.shapes.add_table(num_rows, num_cols, Inches(tbl_left), Inches(tbl_top), Inches(tbl_width), Inches(0.35 * num_rows))
            tbl = tbl_shape.table

            # Header
            headers = ["Metric", "Basic", "With Refinement", "Formula"]
            for ci, h in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = DIM

            # Data rows
            for ri, row_data in enumerate(table_data):
                for ci, val in enumerate(row_data):
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = val
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = WHITE if ci < 2 else DIM
                        if ci == 2 and val in ["0.0012", "58.35 dB", "1.000", "2.3×10⁻⁷"]:
                            p.font.color.rgb = AMBER
                            p.font.bold = True

        # Right panel: What the Metrics Mean
        add_card(slide, 6.6, y_cursor, 5.8, 2.0)
        add_txt(slide, "What the Metrics Mean", 6.9, y_cursor + 0.05, 5.3, 0.4, font_size=15, bold=True, color=WHITE)

        metrics_data = [
            ("RMSE", "Average pixel error. 0.022 = 2.2% per pixel. With refinement: 0.12% — near perfect.", AMBER),
            ("PSNR", "Peak Signal-to-Noise Ratio. >30 dB is good; >50 dB is excellent. Refinement achieves 58.35 dB.", PURPLE),
            ("SSIM", "Structural Similarity — compares patterns, not just pixels. 1.000 = structurally identical.", GREEN),
        ]
        for i, (name, desc, color) in enumerate(metrics_data):
            cy = y_cursor + 0.45 + i * 0.5
            add_txt(slide, name, 7.1, cy, 1.5, 0.3, font_size=12, bold=True, color=color)
            add_txt(slide, desc, 7.1, cy + 0.25, 5.0, 0.3, font_size=9, color=GRAY)

        # Highlight box
        hl_match = re.search(r'<div class="highlight-box">\s*<p>(.*?)</p>', slide_html, re.DOTALL)
        if hl_match:
            hl_text = strip_tags(hl_match.group(1))
            add_txt(slide, hl_text, 0.8, 5.0, 11.5, 0.4, font_size=12, color=GRAY)
        continue

    # ── SLIDE 8: NOISE ROBUSTNESS ──
    if idx == 7:
        # Noise table
        table_data = []
        rows = re.findall(r'<tr>(.*?)</tr>', slide_html, re.DOTALL)
        for row in rows[1:]:
            cells = re.findall(r'<td[^>]*>(.*?)</td>', row, re.DOTALL)
            if cells:
                table_data.append([strip_tags(c) for c in cells])

        if table_data:
            tbl_shape = slide.shapes.add_table(len(table_data) + 1, 4, Inches(0.8), Inches(y_cursor), Inches(4.5), Inches(0.3 * (len(table_data) + 1)))
            tbl = tbl_shape.table
            headers = ["Noise", "RMSE", "PSNR", "SSIM"]
            for ci, h in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = DIM
            for ri, row_data in enumerate(table_data):
                for ci, val in enumerate(row_data):
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = val
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = WHITE

        # Right: Tikhonov card
        add_card(slide, 5.6, y_cursor, 6.8, 3.0)
        add_txt(slide, "Tikhonov Regularization", 5.9, y_cursor + 0.05, 6.3, 0.4, font_size=15, bold=True, color=AMBER)
        add_txt(slide, "Instead of solving min ‖A·x − b‖², we solve:", 5.9, y_cursor + 0.45, 6.3, 0.3, font_size=11, color=GRAY)
        add_txt(slide, "min ‖A·x − b‖² + λ²·‖x‖²", 5.9, y_cursor + 0.8, 6.3, 0.35, font_size=14, color=WHITE, align=PP_ALIGN.CENTER, font_name='Consolas')
        add_txt(slide, "The second term penalizes large solutions, preventing noise amplification. Parameter λ (damp) controls the bias-variance tradeoff.", 5.9, y_cursor + 1.2, 6.3, 0.5, font_size=10, color=GRAY)
        add_txt(slide, "Without regularization: RMSE jumps from 0.022 → 2.56 at just 1% noise! With regularization: RMSE stays at 0.095 — 27× better.", 5.9, y_cursor + 1.8, 6.3, 0.5, font_size=10, color=GRAY)
        continue

    # ── SLIDE 9: SUMMARY ──
    if idx == 8:
        # Core Achievement card
        add_card(slide, 0.8, y_cursor, 5.5, 1.5)
        add_txt(slide, "Core Achievement", 1.1, y_cursor + 0.05, 5.0, 0.35, font_size=14, bold=True, color=AMBER)
        ach_match = re.search(r'Core Achievement</h4>\s*<p>(.*?)</p>', slide_html, re.DOTALL)
        if ach_match:
            add_txt(slide, strip_tags(ach_match.group(1)), 1.1, y_cursor + 0.4, 5.0, 1.0, font_size=11, color=GRAY)

        # What We Learned card
        add_card(slide, 6.6, y_cursor, 5.8, 1.5)
        add_txt(slide, "What We Learned", 6.9, y_cursor + 0.05, 5.3, 0.35, font_size=14, bold=True, color=PURPLE)
        learned_matches = re.findall(r'<p>\d+\.\s*(.*?)</p>', slide_html, re.DOTALL)
        # Get the ones under "What We Learned" - they're at the end
        learned_texts = re.findall(r'<p>(\d+\.\s*<strong[^>]*>.*?</strong>.*?)</p>', slide_html, re.DOTALL)
        if learned_texts:
            for i, lt in enumerate(learned_texts):
                lt_clean = strip_tags(lt)
                add_txt(slide, lt_clean, 6.9, y_cursor + 0.4 + i * 0.25, 5.3, 0.25, font_size=10, color=GRAY)

        y_cursor += 1.8

        # Key Results card
        add_card(slide, 0.8, y_cursor, 5.5, 1.5)
        add_txt(slide, "Key Results", 1.1, y_cursor + 0.05, 5.0, 0.35, font_size=14, bold=True, color=GREEN)
        results = re.findall(r'<p>•\s*(.*?)</p>', slide_html, re.DOTALL)
        if results:
            for i, r in enumerate(results):
                r_clean = strip_tags(r)
                add_txt(slide, f"• {r_clean}", 1.1, y_cursor + 0.4 + i * 0.22, 5.0, 0.22, font_size=10, color=GRAY)

        # Tags at bottom
        tags_match = re.findall(r'<span class="tag[^"]*"[^>]*>([^<]+)</span>', slide_html, re.DOTALL)
        # Get the lowercase tags (tech stack), not the section tags
        tech_tags = ["Python", "NumPy/SciPy", "Matplotlib", "scikit-image", "pydicom", "pytest"]
        tags_str = "  |  ".join(tech_tags)
        add_txt(slide, tags_str, 0.8, 5.8, 11.5, 0.3, font_size=12, color=DIM, align=PP_ALIGN.CENTER)

        # Footer
        add_txt(slide, "CT Reconstruction using LU Decomposition  ·  College Project  ·  Python", 0.8, 6.3, 11.5, 0.3, font_size=10, color=DIM, align=PP_ALIGN.CENTER)
        continue

# Save
output_path = 'presentation.pptx'
prs.save(output_path)
print(f"Saved to {output_path}")
